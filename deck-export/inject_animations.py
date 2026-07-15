# -*- coding: utf-8 -*-
"""
把原生 PowerPoint 动画 (fade 入场 + 错峰依次出现) 注入到现有可编辑 PPTX。
动画用 <p:timing> 写入每页 XML；文字/形状保持可编辑，WPS / PowerPoint 放映时真实播放。
"""
import os
from pptx import Presentation
from pptx.oxml import parse_xml

SRC = 'D:/creditbank-1/deck-export/星秩存册-界面设计答辩-editable.pptx'
OUT = 'D:/creditbank-1/deck-export/星秩存册-界面设计答辩-动画版.pptx'
PNS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

prs = Presentation(SRC)
sw, sh = prs.slide_width, prs.slide_height


def is_background(shape):
    """全屏、无文字的形状视为背景，不对其做动画（避免整页闪白）。"""
    try:
        has_text = shape.has_text_frame and shape.text_frame.text.strip() != ''
    except Exception:
        has_text = False
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return False
    if has_text or None in (l, t, w, h):
        return False
    return (l <= sw * 0.03) and (t <= sh * 0.03) and (w >= sw * 0.95) and (h >= sh * 0.95)


def shape_id(shape):
    cnvpr = shape._element.find('.//{%s}cNvPr' % PNS)
    return int(cnvpr.get('id')) if cnvpr is not None else None


def build_timing(ids, stagger_ms=140, dur_ms=420):
    pars, cid = [], 2
    for i, sid in enumerate(ids):
        delay, bid, aid, animid = i * stagger_ms, cid, cid + 1, cid + 2
        cid += 3
        pars.append(f'''<p:par>
  <p:cTn id="{bid}" fill="hold" dur="1" presetID="1" presetClass="entr" presetSubtype="0" nodeType="withEffect">
    <p:stCondLst><p:cond evt="onBegin" delay="{delay}"/></p:stCondLst>
    <p:childTnLst><p:par>
      <p:cTn id="{aid}" fill="hold" dur="1" nodeType="animEffect">
        <p:stCondLst><p:cond evt="onBegin" delay="0"/></p:stCondLst>
        <p:childTnLst><p:animEffect preset="fade" presetClass="entr" presetSubtype="0" transition="in">
          <p:cBhvr>
            <p:cTn id="{animid}" dur="{dur_ms}" fill="hold" nodeType="anim" presetID="1" presetClass="entr" presetSubtype="0">
              <p:stCondLst><p:cond evt="onBegin" delay="0"/></p:stCondLst>
              <p:to var="style.opacity" val="1"/>
              <p:from var="style.opacity" val="0"/>
            </p:cTn>
            <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
          </p:cBhvr>
        </p:animEffect></p:childTnLst>
      </p:cTn>
    </p:par></p:childTnLst>
    <p:subTnLst/>
  </p:cTn>
</p:par>''')
    after_id = cid
    return f'''<p:timing xmlns:p="{PNS}">
  <p:tnLst><p:par>
    <p:cTn id="1" fill="hold" dur="indefinite" restart="never" nodeType="mainSeq">
      <p:childTnLst>{''.join(pars)}</p:childTnLst>
      <p:subTnLst/>
    </p:cTn>
  </p:par><p:par>
    <p:cTn id="{after_id}" fill="hold" dur="indefinite" nodeType="afterEffectSeq"/>
  </p:par></p:tnLst>
</p:timing>'''


total = 0
for idx, slide in enumerate(prs.slides, 1):
    anim_ids = [shape_id(s) for s in slide.shapes if shape_id(s) is not None and not is_background(s)]
    if not anim_ids:
        continue
    el = parse_xml(build_timing(anim_ids))
    sld = slide._element
    inserted = False
    for i, child in enumerate(sld):
        if child.tag.endswith('}extLst'):
            sld.insert(i, el)
            inserted = True
            break
    if not inserted:
        sld.append(el)
    total += len(anim_ids)
    print(f'slide {idx:>2}: 注入 {len(anim_ids):>2} 个动画')

prs.save(OUT)
print(f'\nOK -> {OUT}')
print(f'总动画数: {total} | 文件大小: {os.path.getsize(OUT)/1024:.0f} KB')
