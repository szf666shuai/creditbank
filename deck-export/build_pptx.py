from pptx import Presentation
from pptx.util import Inches
import os

png_dir = r"D:\creditbank-1\deck-export\png"
out = r"D:\creditbank-1\deck-export\星秩存册-界面设计答辩.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

count = 0
for n in range(2, 30):
    p = os.path.join(png_dir, f"12_{n}.png")
    if not os.path.exists(p):
        print("MISSING", p)
        continue
    slide = prs.slides.add_slide(blank)
    # full-bleed image, exact 16:9 match -> no distortion
    slide.shapes.add_picture(p, 0, 0, prs.slide_width, prs.slide_height)
    count += 1

prs.save(out)
size_mb = os.path.getsize(out) / (1024 * 1024)
print(f"Saved -> {out}")
print(f"Slides: {count} | Size: {size_mb:.1f} MB")
