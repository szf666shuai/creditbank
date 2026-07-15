from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import html
import os

SLIDES_DIR = os.path.join(os.path.dirname(__file__), 'slides')
OUTPUT_FILE = os.path.join(os.path.dirname(__file__), '学分银行平台答辩PPT.pptx')

COLORS = {
    'ink': RGBColor(26, 32, 44),
    'cream': RGBColor(255, 249, 240),
    'green': RGBColor(34, 197, 94),
    'green_deep': RGBColor(22, 163, 74),
    'blue': RGBColor(190, 227, 248),
    'pink': RGBColor(254, 205, 211),
    'yellow': RGBColor(254, 240, 138),
    'purple': RGBColor(221, 214, 254),
    'white': RGBColor(255, 255, 255),
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    slide_files = sorted([f for f in os.listdir(SLIDES_DIR) if f.endswith('.html')])
    
    for idx, slide_file in enumerate(slide_files, 1):
        slide_path = os.path.join(SLIDES_DIR, slide_file)
        with open(slide_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        tree = html.fromstring(content)
        page_number = tree.xpath('//div[@class="page-number"]/text()')
        page_title = tree.xpath('//h2/text()')
        
        page_num = page_number[0] if page_number else f'{idx:02d} / {len(slide_files)}'
        title_text = page_title[0] if page_title else 'Slide'
        
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['cream']
        
        add_page_header(slide, page_num, title_text)
        
        if slide_file == '01-cover.html':
            add_cover_content(slide)
        elif slide_file == '02-agenda.html':
            add_agenda_content(slide, tree)
        elif slide_file == '03-story.html':
            add_story_content(slide, tree)
        elif slide_file == '04-positioning.html':
            add_positioning_content(slide, tree)
        elif slide_file == '05-color.html':
            add_color_content(slide, tree)
        elif slide_file == '06-functions.html':
            add_functions_content(slide, tree)
        elif slide_file == '07-agent.html':
            add_agent_content(slide, tree)
        elif slide_file == '08-tech.html':
            add_tech_content(slide, tree)
        elif slide_file == '09-innovation.html':
            add_innovation_content(slide, tree)
        elif slide_file == '10-process.html':
            add_process_content(slide, tree)
        elif slide_file == '11-team.html':
            add_team_content(slide, tree)
        elif slide_file == '12-summary.html':
            add_summary_content(slide, tree)
    
    prs.save(OUTPUT_FILE)
    print(f'PPT saved to: {OUTPUT_FILE}')

def add_page_header(slide, page_number, title):
    txBox = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(10), Inches(1))
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = page_number
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['green']
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True

def add_cover_content(slide):
    title_box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(10), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = '学分银行平台'
    p.font.size = Pt(56)
    p.font.color.rgb = COLORS['green']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(4), Inches(5), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = '终身学习 · 成就未来'
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['ink']
    p.alignment = PP_ALIGN.CENTER
    
    badge_box = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(4), Inches(1))
    tf = badge_box.text_frame
    p = tf.add_paragraph()
    p.text = '学习的健身房'
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['green']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    green_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(2))
    green_rect.fill.solid()
    green_rect.fill.fore_color.rgb = COLORS['green']
    green_rect.line.color.rgb = COLORS['ink']
    green_rect.line.width = Pt(3)
    
    blue_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13), Inches(6), Inches(2), Inches(2))
    blue_rect.fill.solid()
    blue_rect.fill.fore_color.rgb = COLORS['blue']
    blue_rect.line.color.rgb = COLORS['ink']
    blue_rect.line.width = Pt(3)

def add_agenda_content(slide, tree):
    cards = tree.xpath('//div[@class="agenda-card"]')
    rows = 2
    cols = 3
    card_width = Inches(2.2)
    card_height = Inches(2.5)
    start_x = Inches(1.5)
    start_y = Inches(2)
    gap_x = Inches(0.8)
    gap_y = Inches(1)
    
    color_map = [COLORS['green'], COLORS['blue'], COLORS['pink'], COLORS['yellow'], COLORS['purple'], COLORS['blue']]
    
    for i, card in enumerate(cards):
        icon = card.xpath('.//div[@class="agenda-card__icon"]/text()')
        title = card.xpath('.//div[@class="agenda-card__title"]/text()')
        desc = card.xpath('.//div[@class="agenda-card__desc"]/text()')
        
        row = i // cols
        col = i % cols
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color_map[i % len(color_map)]
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        content_x = x + Inches(0.2)
        content_y = y + Inches(0.2)
        
        icon_box = slide.shapes.add_textbox(content_x, content_y, Inches(0.8), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = icon[0] if icon else '📌'
        ip.font.size = Pt(24)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(content_x, content_y + Inches(0.8), card_width - Inches(0.4), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = title[0] if title else ''
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        desc_box = slide.shapes.add_textbox(content_x, content_y + Inches(1.3), card_width - Inches(0.4), Inches(0.8))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = desc[0] if desc else ''
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['ink']

def add_story_content(slide, tree):
    pain_points = tree.xpath('//div[@class="pain-point"]')
    card_width = Inches(5)
    card_height = Inches(2.5)
    start_x = Inches(0.8)
    start_y = Inches(2)
    gap_x = Inches(0.8)
    
    for i, point in enumerate(pain_points):
        title = point.xpath('.//div[@class="pain-point__title"]/text()')
        desc = point.xpath('.//div[@class="pain-point__desc"]/text()')
        
        x = start_x + i * (card_width + gap_x)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['pink'] if i == 0 else (COLORS['yellow'] if i == 1 else COLORS['blue'])
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        title_box = slide.shapes.add_textbox(x + Inches(0.3), start_y + Inches(0.3), card_width - Inches(0.6), Inches(0.5))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = title[0] if title else ''
        tp.font.size = Pt(16)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), start_y + Inches(0.9), card_width - Inches(0.6), Inches(1.2))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = desc[0] if desc else ''
        dp.font.size = Pt(11)
        dp.font.color.rgb = COLORS['ink']
    
    solution_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(13), Inches(2))
    sb = solution_box.text_frame
    sp = sb.add_paragraph()
    sp.text = '学分银行平台 - 连接学习与就业，打造终身学习生态圈'
    sp.font.size = Pt(22)
    sp.font.color.rgb = COLORS['green']
    sp.font.bold = True

def add_positioning_content(slide, tree):
    left_x = Inches(1)
    right_x = Inches(9)
    card_width = Inches(7)
    card_height = Inches(3.5)
    start_y = Inches(2)
    
    left_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, start_y, card_width, card_height)
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = COLORS['green']
    left_rect.line.color.rgb = COLORS['ink']
    left_rect.line.width = Pt(3)
    
    icon_box = slide.shapes.add_textbox(left_x + Inches(3), start_y + Inches(0.5), Inches(1), Inches(1))
    ib = icon_box.text_frame
    ip = ib.add_paragraph()
    ip.text = '🏋️'
    ip.font.size = Pt(36)
    ip.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(left_x + Inches(1), start_y + Inches(1.6), card_width - Inches(2), Inches(0.5))
    tb = title_box.text_frame
    tp = tb.add_paragraph()
    tp.text = '学习的健身房'
    tp.font.size = Pt(24)
    tp.font.color.rgb = COLORS['white']
    tp.font.bold = True
    tp.alignment = PP_ALIGN.CENTER
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, start_y, card_width, card_height)
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['blue']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    icon_box2 = slide.shapes.add_textbox(right_x + Inches(3), start_y + Inches(0.5), Inches(1), Inches(1))
    ib2 = icon_box2.text_frame
    ip2 = ib2.add_paragraph()
    ip2.text = '📚'
    ip2.font.size = Pt(36)
    ip2.alignment = PP_ALIGN.CENTER
    
    title_box2 = slide.shapes.add_textbox(right_x + Inches(1), start_y + Inches(1.6), card_width - Inches(2), Inches(0.5))
    tb2 = title_box2.text_frame
    tp2 = tb2.add_paragraph()
    tp2.text = '知识的图书馆'
    tp2.font.size = Pt(24)
    tp2.font.color.rgb = COLORS['ink']
    tp2.font.bold = True
    tp2.alignment = PP_ALIGN.CENTER
    
    desc_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(12), Inches(1.5))
    db = desc_box.text_frame
    dp = db.add_paragraph()
    dp.text = '为什么在家能学习还要去图书馆？为什么在家能锻炼还要去健身房？因为氛围与心理暗示无可替代。我们的平台就是学习领域的健身房与图书馆。'
    dp.font.size = Pt(14)
    dp.font.color.rgb = COLORS['ink']
    dp.alignment = PP_ALIGN.CENTER

def add_color_content(slide, tree):
    color_items = [
        {'color': COLORS['green'], 'name': '绿色', 'desc': '成长、希望、活力、教育'},
        {'color': COLORS['cream'], 'name': '奶油色', 'desc': '温暖、舒适、包容、亲和力'},
        {'color': COLORS['blue'], 'name': '蓝色', 'desc': '专业、信任、科技感'},
        {'color': COLORS['yellow'], 'name': '黄色', 'desc': '活力、创意、注意力'},
    ]
    
    card_width = Inches(3.5)
    card_height = Inches(2)
    start_x = Inches(1)
    start_y = Inches(2)
    gap_x = Inches(0.8)
    
    for i, item in enumerate(color_items):
        x = start_x + i * (card_width + gap_x)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = item['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        title_box = slide.shapes.add_textbox(x + Inches(0.3), start_y + Inches(0.3), card_width - Inches(0.6), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = item['name']
        tp.font.size = Pt(16)
        tp.font.color.rgb = COLORS['ink'] if item['color'] != COLORS['green'] else COLORS['white']
        tp.font.bold = True
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), start_y + Inches(0.8), card_width - Inches(0.6), Inches(0.8))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = item['desc']
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['ink'] if item['color'] != COLORS['green'] else COLORS['white']
    
    design_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(2))
    db = design_box.text_frame
    dp = db.add_paragraph()
    dp.text = 'Neo-brutal设计风格：粗边框、硬阴影、大胆配色，打造独特的视觉识别度，吸引各年龄段用户'
    dp.font.size = Pt(16)
    dp.font.color.rgb = COLORS['ink']
    dp.font.bold = True

def add_functions_content(slide, tree):
    cards = tree.xpath('//div[@class="function-card"]')
    rows = 2
    cols = 3
    card_width = Inches(2.6)
    card_height = Inches(3)
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    gap_x = Inches(0.5)
    gap_y = Inches(0.5)
    
    color_map = [COLORS['green'], COLORS['blue'], COLORS['pink'], COLORS['yellow'], COLORS['purple'], COLORS['blue']]
    
    for i, card in enumerate(cards):
        icon = card.xpath('.//div[@class="function-card__icon"]/text()')
        title = card.xpath('.//div[@class="function-card__title"]/text()')
        
        row = i // cols
        col = i % cols
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['green'] if i == 0 else color_map[i % len(color_map)]
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.3), Inches(1), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = icon[0] if icon else '📌'
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.2), card_width - Inches(0.4), Inches(0.5))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = title[0] if title else ''
        tp.font.size = Pt(13)
        tp.font.color.rgb = COLORS['white'] if i == 0 else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER

def add_agent_content(slide, tree):
    left_x = Inches(1)
    right_x = Inches(10)
    card_width = Inches(8)
    card_height = Inches(3.5)
    start_y = Inches(2)
    
    left_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, start_y, card_width, card_height)
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = COLORS['green']
    left_rect.line.color.rgb = COLORS['ink']
    left_rect.line.width = Pt(3)
    
    icon_box = slide.shapes.add_textbox(left_x + Inches(3.5), start_y + Inches(0.5), Inches(1), Inches(1))
    ib = icon_box.text_frame
    ip = ib.add_paragraph()
    ip.text = '🤖'
    ip.font.size = Pt(40)
    ip.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(left_x + Inches(1), start_y + Inches(1.6), card_width - Inches(2), Inches(0.5))
    tb = title_box.text_frame
    tp = tb.add_paragraph()
    tp.text = '学习伙伴'
    tp.font.size = Pt(24)
    tp.font.color.rgb = COLORS['white']
    tp.font.bold = True
    tp.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(left_x + Inches(1), start_y + Inches(2.3), card_width - Inches(2), Inches(0.5))
    sb = subtitle_box.text_frame
    sp = sb.add_paragraph()
    sp.text = '24小时陪伴式学习助手'
    sp.font.size = Pt(14)
    sp.font.color.rgb = COLORS['white']
    sp.alignment = PP_ALIGN.CENTER
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, start_y, card_width, card_height)
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['white']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    chat_box = slide.shapes.add_textbox(right_x + Inches(0.3), start_y + Inches(0.3), card_width - Inches(0.6), Inches(2.5))
    cb = chat_box.text_frame
    
    messages = [
        ('AI', '你好！我是你的学习助手，请问有什么可以帮你的？'),
        ('User', '我想学习Python，有什么推荐的课程吗？'),
        ('AI', '根据你的学习目标，我推荐Python基础入门、数据分析实战等课程。'),
    ]
    
    for sender, text in messages:
        p = cb.add_paragraph()
        p.text = f'{sender}: {text}'
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['green'] if sender == 'AI' else COLORS['blue']
        p.font.bold = True
    
    features = ['智能课程推荐', '实时答疑', '简历生成', '学习规划']
    feature_width = Inches(1.8)
    feature_height = Inches(0.8)
    feature_start_x = Inches(1.5)
    feature_start_y = Inches(6)
    feature_gap = Inches(0.5)
    
    for i, feature in enumerate(features):
        x = feature_start_x + i * (feature_width + feature_gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, feature_start_y, feature_width, feature_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['white']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(2)
        
        text_box = slide.shapes.add_textbox(x + Inches(0.1), feature_start_y + Inches(0.1), feature_width - Inches(0.2), Inches(0.6))
        tb = text_box.text_frame
        tp = tb.add_paragraph()
        tp.text = feature
        tp.font.size = Pt(11)
        tp.font.color.rgb = COLORS['ink']
        tp.alignment = PP_ALIGN.CENTER

def add_tech_content(slide, tree):
    left_x = Inches(1)
    right_x = Inches(8.5)
    card_width = Inches(7)
    card_height = Inches(2.5)
    start_y = Inches(1.8)
    
    layers = [
        {'name': '前端层', 'techs': ['Vue 3', 'Element Plus', 'Vite', 'Pinia'], 'color': COLORS['pink']},
        {'name': '后端层', 'techs': ['Spring Boot', 'MyBatis Plus', 'JWT', 'WebSocket'], 'color': COLORS['green']},
        {'name': '数据层', 'techs': ['MySQL 8.0', 'Redis 7'], 'color': COLORS['yellow']},
    ]
    
    for i, layer in enumerate(layers):
        y = start_y + i * (card_height + Inches(0.3))
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = layer['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        title_box = slide.shapes.add_textbox(left_x + Inches(0.3), y + Inches(0.2), Inches(2), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = layer['name']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        tech_box = slide.shapes.add_textbox(left_x + Inches(0.3), y + Inches(0.7), card_width - Inches(0.6), Inches(1.4))
        tb = tech_box.text_frame
        for tech in layer['techs']:
            p = tb.add_paragraph()
            p.text = f'• {tech}'
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['ink']
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, start_y, card_width, card_height * 3 + Inches(0.6))
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['blue']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    stacks = [
        {'title': '前端技术栈', 'items': ['Vue 3.5.x', 'Vite 8.1.x', 'Element Plus 2.14.x', 'TypeScript 6.0.x']},
        {'title': '后端技术栈', 'items': ['Spring Boot 3.4.x', 'MyBatis Plus 3.5.x', 'MySQL 8.0', 'Redis 7.0']},
        {'title': '部署架构', 'items': ['Docker 24.x', 'Docker Compose 2.x', 'Nginx 1.25.x', 'TRTC 5.18.x']},
    ]
    
    stack_y = start_y + Inches(0.3)
    for stack in stacks:
        title_box = slide.shapes.add_textbox(right_x + Inches(0.3), stack_y, card_width - Inches(0.6), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = stack['title']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        item_box = slide.shapes.add_textbox(right_x + Inches(0.3), stack_y + Inches(0.5), card_width - Inches(0.6), Inches(1))
        ib = item_box.text_frame
        for item in stack['items']:
            p = ib.add_paragraph()
            p.text = f'• {item}'
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['ink']
        
        stack_y += Inches(1.8)

def add_innovation_content(slide, tree):
    innovations = [
        {'icon': '🤖', 'title': 'AI驱动学习', 'desc': '深度整合大语言模型，个性化学习推荐', 'color': COLORS['green']},
        {'icon': '💰', 'title': '学分银行体系', 'desc': '学习成果量化，建立完整学分流通体系', 'color': COLORS['yellow']},
        {'icon': '📹', 'title': '实时视频面试', 'desc': '整合TRTC，实现在线视频面试', 'color': COLORS['blue']},
        {'icon': '🎨', 'title': 'Neo-brutal风格', 'desc': '独特视觉风格，提升品牌识别度', 'color': COLORS['pink']},
        {'icon': '⚡', 'title': '前后端分离', 'desc': 'Vue3+Spring Boot，快速迭代部署', 'color': COLORS['purple']},
    ]
    
    card_width = Inches(3.4)
    card_height = Inches(2.8)
    start_x = Inches(0.8)
    start_y = Inches(2)
    gap_x = Inches(0.5)
    gap_y = Inches(0.5)
    
    for i, item in enumerate(innovations):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = item['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.3), Inches(1), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = item['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.2), card_width - Inches(0.4), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = item['title']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.7), card_width - Inches(0.4), Inches(0.8))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = item['desc']
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['ink']
        dp.alignment = PP_ALIGN.CENTER

def add_process_content(slide, tree):
    processes = [
        {'title': '面试流程', 'icon': '🎤', 'steps': ['浏览职位', '投递简历', '面试邀请', '在线面试', '面试结果'], 'color': COLORS['blue']},
        {'title': '活动流程', 'icon': '🎉', 'steps': ['发布活动', '邀请学员', '报名参加', '活动签到', '获取学分'], 'color': COLORS['pink']},
        {'title': '学习流程', 'icon': '📚', 'steps': ['选择课程', '开始学习', '互动交流', '完成课程', '积累学分'], 'color': COLORS['green']},
    ]
    
    card_width = Inches(5)
    card_height = Inches(5)
    start_x = Inches(0.5)
    start_y = Inches(2)
    gap_x = Inches(0.6)
    
    for i, process in enumerate(processes):
        x = start_x + i * (card_width + gap_x)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = process['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(2), start_y + Inches(0.3), Inches(1), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = process['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(1), start_y + Inches(1.2), card_width - Inches(2), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = process['title']
        tp.font.size = Pt(18)
        tp.font.color.rgb = COLORS['white'] if process['color'] == COLORS['green'] else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        step_y = start_y + Inches(1.8)
        for j, step in enumerate(process['steps']):
            step_box = slide.shapes.add_textbox(x + Inches(0.3), step_y, card_width - Inches(0.6), Inches(0.6))
            sb = step_box.text_frame
            sp = sb.add_paragraph()
            sp.text = f'{j+1:02d}. {step}'
            sp.font.size = Pt(12)
            sp.font.color.rgb = COLORS['white'] if process['color'] == COLORS['green'] else COLORS['ink']
            step_y += Inches(0.7)

def add_team_content(slide, tree):
    members = [
        {'icon': '👨‍💻', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]'},
        {'icon': '👩‍💻', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]'},
        {'icon': '👨‍🎨', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]'},
        {'icon': '👩‍🎨', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]'},
    ]
    
    card_width = Inches(3.8)
    card_height = Inches(2.2)
    start_x = Inches(0.8)
    start_y = Inches(2)
    gap_x = Inches(0.6)
    
    color_map = [COLORS['blue'], COLORS['pink'], COLORS['yellow'], COLORS['purple']]
    
    for i, member in enumerate(members):
        x = start_x + i * (card_width + gap_x)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color_map[i]
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(0.3), start_y + Inches(0.2), Inches(1), Inches(1))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = member['icon']
        ip.font.size = Pt(32)
        ip.alignment = PP_ALIGN.CENTER
        
        info_box = slide.shapes.add_textbox(x + Inches(1.5), start_y + Inches(0.2), card_width - Inches(2), Inches(1.8))
        ib = info_box.text_frame
        p = ib.add_paragraph()
        p.text = member['name']
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['ink']
        p.font.bold = True
        
        p = ib.add_paragraph()
        p.text = member['role']
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['ink']
        
        p = ib.add_paragraph()
        p.text = member['desc']
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['ink']
    
    workflow_y = Inches(4.8)
    workflow_steps = [
        ('01', '需求分析', '收集用户需求，分析业务场景'),
        ('02', '系统设计', '架构设计、数据库设计'),
        ('03', '开发实现', '前后端并行开发'),
        ('04', '部署上线', 'Docker容器化部署'),
    ]
    
    step_width = Inches(3.8)
    step_height = Inches(1.2)
    
    for i, (num, title, desc) in enumerate(workflow_steps):
        x = start_x + i * (step_width + gap_x)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, workflow_y, step_width, step_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['white']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(2)
        
        num_box = slide.shapes.add_textbox(x + Inches(0.2), workflow_y + Inches(0.1), Inches(0.5), Inches(0.8))
        nb = num_box.text_frame
        np = nb.add_paragraph()
        np.text = num
        np.font.size = Pt(20)
        np.font.color.rgb = COLORS['green']
        np.font.bold = True
        
        title_box = slide.shapes.add_textbox(x + Inches(0.8), workflow_y + Inches(0.1), Inches(1.5), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = title
        tp.font.size = Pt(12)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.8), workflow_y + Inches(0.5), step_width - Inches(1), Inches(0.5))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = desc
        dp.font.size = Pt(9)
        dp.font.color.rgb = COLORS['ink']

def add_summary_content(slide, tree):
    achievements = [
        ('完整的学习平台', '课程学习、学分银行、社区论坛'),
        ('AI智能助手', '个性化学习推荐和智能答疑'),
        ('企业对接功能', '企业入驻、职位发布、在线面试'),
        ('Neo-brutal设计', '独特视觉风格，提升品牌识别度'),
    ]
    
    future_items = [
        '深化AI应用 - AI课程生成、智能评估',
        '扩展社交功能 - 学习小组、好友互动',
        '移动端适配 - 小程序和APP',
        '数据分析平台 - 可视化报表和洞察',
    ]
    
    left_x = Inches(1)
    right_x = Inches(9)
    card_width = Inches(7)
    card_height = Inches(3.5)
    start_y = Inches(2)
    
    left_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, start_y, card_width, card_height)
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = COLORS['white']
    left_rect.line.color.rgb = COLORS['ink']
    left_rect.line.width = Pt(3)
    
    title_box = slide.shapes.add_textbox(left_x + Inches(0.3), start_y + Inches(0.2), card_width - Inches(0.6), Inches(0.4))
    tb = title_box.text_frame
    tp = tb.add_paragraph()
    tp.text = '项目成果'
    tp.font.size = Pt(18)
    tp.font.color.rgb = COLORS['ink']
    tp.font.bold = True
    
    item_y = start_y + Inches(0.8)
    for title, desc in achievements:
        item_box = slide.shapes.add_textbox(left_x + Inches(0.3), item_y, card_width - Inches(0.6), Inches(0.7))
        ib = item_box.text_frame
        p = ib.add_paragraph()
        p.text = f'✅ {title}'
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['ink']
        p.font.bold = True
        
        p = ib.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['ink']
        item_y += Inches(0.8)
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, start_y, card_width, card_height)
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['white']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    title_box2 = slide.shapes.add_textbox(right_x + Inches(0.3), start_y + Inches(0.2), card_width - Inches(0.6), Inches(0.4))
    tb2 = title_box2.text_frame
    tp2 = tb2.add_paragraph()
    tp2.text = '未来展望'
    tp2.font.size = Pt(18)
    tp2.font.color.rgb = COLORS['ink']
    tp2.font.bold = True
    
    item_y2 = start_y + Inches(0.8)
    for item in future_items:
        item_box2 = slide.shapes.add_textbox(right_x + Inches(0.3), item_y2, card_width - Inches(0.6), Inches(0.6))
        ib2 = item_box2.text_frame
        p = ib2.add_paragraph()
        p.text = f'💡 {item}'
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['ink']
        item_y2 += Inches(0.7)
    
    numbers = [('6', '核心模块'), ('5', '技术创新'), ('3', '核心流程'), ('4', '团队成员')]
    num_width = Inches(3.5)
    num_height = Inches(1)
    num_start_x = Inches(1)
    num_start_y = Inches(6.2)
    num_gap = Inches(0.5)
    
    for i, (num, label) in enumerate(numbers):
        x = num_start_x + i * (num_width + num_gap)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, num_start_y, num_width, num_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['white']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(2)
        
        num_box = slide.shapes.add_textbox(x + Inches(0.3), num_start_y + Inches(0.1), Inches(1), Inches(0.8))
        nb = num_box.text_frame
        np = nb.add_paragraph()
        np.text = num
        np.font.size = Pt(32)
        np.font.color.rgb = COLORS['green']
        np.font.bold = True
        
        label_box = slide.shapes.add_textbox(x + Inches(1.5), num_start_y + Inches(0.2), Inches(1.8), Inches(0.6))
        lb = label_box.text_frame
        lp = lb.add_paragraph()
        lp.text = label
        lp.font.size = Pt(14)
        lp.font.color.rgb = COLORS['ink']
        lp.font.bold = True
    
    closing_box = slide.shapes.add_textbox(Inches(6), Inches(7.8), Inches(4), Inches(1))
    cb = closing_box.text_frame
    cp = cb.add_paragraph()
    cp.text = '感谢聆听！'
    cp.font.size = Pt(28)
    cp.font.color.rgb = COLORS['green']
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER

if __name__ == '__main__':
    create_presentation()