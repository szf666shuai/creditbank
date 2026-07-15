from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_FILE = os.path.join(os.path.dirname(__file__), '学分银行平台答辩PPT_完整版.pptx')

COLORS = {
    'ink': RGBColor(26, 32, 44),
    'cream': RGBColor(255, 249, 240),
    'green': RGBColor(34, 197, 94),
    'green_deep': RGBColor(22, 163, 74),
    'blue': RGBColor(190, 227, 248),
    'pink': RGBColor(254, 205, 211),
    'yellow': RGBColor(254, 240, 138),
    'purple': RGBColor(221, 214, 254),
    'white': RGBColor(255, 255, 255),
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    add_cover(prs)
    add_agenda(prs)
    add_story(prs)
    add_positioning(prs)
    add_color(prs)
    add_functions(prs)
    add_agent(prs)
    add_tech(prs)
    add_innovation(prs)
    add_process(prs)
    add_team(prs)
    add_summary(prs)
    
    prs.save(OUTPUT_FILE)
    print(f'PPT saved to: {OUTPUT_FILE}')

def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    logo_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(2), Inches(2))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = COLORS['green']
    logo_box.line.color.rgb = COLORS['ink']
    logo_box.line.width = Pt(4)
    
    title_box = slide.shapes.add_textbox(Inches(4), Inches(4.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = '学分银行平台'
    p.font.size = Pt(48)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(5), Inches(5.5), Inches(6), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = '终身学习 · 成就未来'
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['ink']
    p.alignment = PP_ALIGN.CENTER
    
    tagline_box = slide.shapes.add_textbox(Inches(5.5), Inches(6.5), Inches(5), Inches(0.8))
    tf = tagline_box.text_frame
    p = tf.add_paragraph()
    p.text = '「让学习成为一种生活方式」'
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['green']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def add_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '02 / 12  答辩大纲'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    cards = [
        {'icon': '📖', 'title': '品牌故事', 'desc': '为什么做这个项目', 'color': COLORS['green']},
        {'icon': '🎨', 'title': '产品设计', 'desc': '定位与视觉设计', 'color': COLORS['blue']},
        {'icon': '⚡', 'title': '核心功能', 'desc': 'AI助手与学分体系', 'color': COLORS['pink']},
        {'icon': '🛠️', 'title': '技术架构', 'desc': '前后端技术栈', 'color': COLORS['yellow']},
        {'icon': '💡', 'title': '创新亮点', 'desc': '差异化竞争力', 'color': COLORS['purple']},
        {'icon': '👥', 'title': '团队总结', 'desc': '成员与成果', 'color': COLORS['blue']},
    ]
    
    for i, card in enumerate(cards):
        row = i // 3
        col = i % 3
        x = Inches(1 + col * 5)
        y = Inches(1.5 + row * 2.8)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(2.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = card['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(1.7), y + Inches(0.3), Inches(1), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = card['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.1), Inches(4), Inches(0.5))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = card['title']
        tp.font.size = Pt(16)
        tp.font.color.rgb = COLORS['white'] if card['color'] == COLORS['green'] else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.7), Inches(4), Inches(0.6))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = card['desc']
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['white'] if card['color'] == COLORS['green'] else COLORS['ink']
        dp.alignment = PP_ALIGN.CENTER

def add_story(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '03 / 12  品牌故事'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    pains = [
        {'icon': '🎓', 'title': '大学与社会的脱节', 'desc': '大学教授传授的知识与社会工作实际需要的技能存在差距', 'color': COLORS['pink']},
        {'icon': '📚', 'title': '终身学习的需求', 'desc': '随着技术快速发展，知识更新周期缩短，人们需要不断学习', 'color': COLORS['yellow']},
        {'icon': '🏢', 'title': '企业人才培养困境', 'desc': '企业难以找到符合岗位要求的人才，培训成本高', 'color': COLORS['blue']},
    ]
    
    for i, pain in enumerate(pains):
        x = Inches(1 + i * 5)
        y = Inches(1.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(3.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = pain['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(1.7), y + Inches(0.3), Inches(1), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = pain['icon']
        ip.font.size = Pt(32)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.2), Inches(4), Inches(0.5))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = pain['title']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.8), Inches(4), Inches(1.4))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = pain['desc']
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['ink']
        dp.alignment = PP_ALIGN.CENTER
    
    solution_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(12), Inches(0.8))
    sb = solution_box.text_frame
    sp = sb.add_paragraph()
    sp.text = '学分银行平台 - 连接学习与就业，打造终身学习生态圈'
    sp.font.size = Pt(20)
    sp.font.color.rgb = COLORS['green']
    sp.font.bold = True
    sp.alignment = PP_ALIGN.CENTER

def add_positioning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '04 / 12  产品定位'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    left_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(6.5), Inches(4))
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = COLORS['green']
    left_rect.line.color.rgb = COLORS['ink']
    left_rect.line.width = Pt(3)
    
    icon_box1 = slide.shapes.add_textbox(Inches(3.7), Inches(2), Inches(1), Inches(1))
    ib = icon_box1.text_frame
    ip = ib.add_paragraph()
    ip.text = '🏋️'
    ip.font.size = Pt(48)
    ip.alignment = PP_ALIGN.CENTER
    
    title_box1 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(5.5), Inches(0.6))
    tb = title_box1.text_frame
    tp = tb.add_paragraph()
    tp.text = '学习的健身房'
    tp.font.size = Pt(24)
    tp.font.color.rgb = COLORS['white']
    tp.font.bold = True
    tp.alignment = PP_ALIGN.CENTER
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.5), Inches(6.5), Inches(4))
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['blue']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    icon_box2 = slide.shapes.add_textbox(Inches(11.2), Inches(2), Inches(1), Inches(1))
    ib = icon_box2.text_frame
    ip = ib.add_paragraph()
    ip.text = '📚'
    ip.font.size = Pt(48)
    ip.alignment = PP_ALIGN.CENTER
    
    title_box2 = slide.shapes.add_textbox(Inches(9), Inches(3.2), Inches(5.5), Inches(0.6))
    tb = title_box2.text_frame
    tp = tb.add_paragraph()
    tp.text = '知识的图书馆'
    tp.font.size = Pt(24)
    tp.font.color.rgb = COLORS['ink']
    tp.font.bold = True
    tp.alignment = PP_ALIGN.CENTER
    
    desc_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(12), Inches(1))
    db = desc_box.text_frame
    dp = db.add_paragraph()
    dp.text = '为什么在家能学习还要去图书馆？因为氛围与心理暗示无可替代'
    dp.font.size = Pt(14)
    dp.font.color.rgb = COLORS['ink']
    dp.alignment = PP_ALIGN.CENTER

def add_color(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '05 / 12  色彩心理学'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    colors = [
        {'name': '绿色', 'desc': '成长、希望、活力、教育', 'color': COLORS['green']},
        {'name': '奶油色', 'desc': '温暖、舒适、包容、亲和力', 'color': COLORS['cream']},
        {'name': '蓝色', 'desc': '专业、信任、科技感', 'color': COLORS['blue']},
        {'name': '黄色', 'desc': '活力、创意、注意力', 'color': COLORS['yellow']},
    ]
    
    for i, item in enumerate(colors):
        x = Inches(1 + i * 3.8)
        y = Inches(1.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.5), Inches(2.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = item['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        title_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(0.4), Inches(2.5), Inches(0.5))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = item['name']
        tp.font.size = Pt(18)
        tp.font.color.rgb = COLORS['white'] if item['color'] == COLORS['green'] else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1), Inches(3), Inches(1.2))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = item['desc']
        dp.font.size = Pt(11)
        dp.font.color.rgb = COLORS['white'] if item['color'] == COLORS['green'] else COLORS['ink']
        dp.alignment = PP_ALIGN.CENTER
    
    design_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(12), Inches(1))
    db = design_box.text_frame
    dp = db.add_paragraph()
    dp.text = 'Neo-brutal设计风格：粗边框、硬阴影、大胆配色，打造独特的视觉识别度'
    dp.font.size = Pt(16)
    dp.font.color.rgb = COLORS['ink']
    dp.font.bold = True
    dp.alignment = PP_ALIGN.CENTER

def add_functions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '06 / 12  核心功能'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    functions = [
        {'icon': '🤖', 'title': '智能学习助手', 'color': COLORS['green']},
        {'icon': '📚', 'title': '学习管理', 'color': COLORS['pink']},
        {'icon': '🏢', 'title': '企业对接', 'color': COLORS['purple']},
        {'icon': '💰', 'title': '学分银行', 'color': COLORS['yellow']},
        {'icon': '💬', 'title': '社区论坛', 'color': COLORS['blue']},
        {'icon': '🎯', 'title': '诚信体系', 'color': COLORS['blue']},
    ]
    
    for i, func in enumerate(functions):
        row = i // 3
        col = i % 3
        x = Inches(1 + col * 5)
        y = Inches(1.5 + row * 2.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(2.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = func['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(1.7), y + Inches(0.3), Inches(1), Inches(0.6))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = func['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1), Inches(3.5), Inches(0.5))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = func['title']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['white'] if func['color'] == COLORS['green'] else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
    
    footer_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(8), Inches(0.6))
    fb = footer_box.text_frame
    fp = fb.add_paragraph()
    fp.text = '6大核心模块，构建完整的终身学习生态'
    fp.font.size = Pt(16)
    fp.font.color.rgb = COLORS['green']
    fp.font.bold = True
    fp.alignment = PP_ALIGN.CENTER

def add_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '07 / 12  AI智能助手'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    left_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(6.5), Inches(4))
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = COLORS['green']
    left_rect.line.color.rgb = COLORS['ink']
    left_rect.line.width = Pt(3)
    
    icon_box1 = slide.shapes.add_textbox(Inches(3.7), Inches(2), Inches(1), Inches(1))
    ib = icon_box1.text_frame
    ip = ib.add_paragraph()
    ip.text = '🤖'
    ip.font.size = Pt(48)
    ip.alignment = PP_ALIGN.CENTER
    
    title_box1 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(5.5), Inches(0.5))
    tb = title_box1.text_frame
    tp = tb.add_paragraph()
    tp.text = '学习伙伴'
    tp.font.size = Pt(24)
    tp.font.color.rgb = COLORS['white']
    tp.font.bold = True
    tp.alignment = PP_ALIGN.CENTER
    
    subtitle_box1 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(5.5), Inches(0.5))
    sb = subtitle_box1.text_frame
    sp = sb.add_paragraph()
    sp.text = '24小时陪伴式学习助手'
    sp.font.size = Pt(14)
    sp.font.color.rgb = COLORS['white']
    sp.alignment = PP_ALIGN.CENTER
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.5), Inches(6.5), Inches(4))
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['white']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    chat_box = slide.shapes.add_textbox(Inches(9), Inches(2), Inches(5.5), Inches(3))
    cb = chat_box.text_frame
    
    messages = [
        ('AI:', '你好！我是你的学习助手，请问有什么可以帮你的？'),
        ('你:', '我想学习Python，有什么推荐的课程吗？'),
        ('AI:', '根据你的学习目标，我推荐Python基础入门等课程。'),
    ]
    
    for sender, text in messages:
        p = cb.add_paragraph()
        p.text = f'{sender} {text}'
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['green'] if sender.startswith('AI') else COLORS['blue']
        p.font.bold = True
    
    features = ['智能课程推荐', '实时答疑', '简历生成', '学习规划']
    for i, feature in enumerate(features):
        x = Inches(1 + i * 3.8)
        y = Inches(6)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.5), Inches(0.8))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['white']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(2)
        
        text_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.1), Inches(3), Inches(0.6))
        tb = text_box.text_frame
        tp = tb.add_paragraph()
        tp.text = feature
        tp.font.size = Pt(12)
        tp.font.color.rgb = COLORS['ink']
        tp.alignment = PP_ALIGN.CENTER

def add_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '08 / 12  技术架构'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    layers = [
        {'name': '前端层', 'techs': ['Vue 3', 'Element Plus', 'Vite', 'Pinia'], 'color': COLORS['pink']},
        {'name': '后端层', 'techs': ['Spring Boot', 'MyBatis Plus', 'JWT', 'WebSocket'], 'color': COLORS['green']},
        {'name': '数据层', 'techs': ['MySQL 8.0', 'Redis 7'], 'color': COLORS['yellow']},
    ]
    
    for i, layer in enumerate(layers):
        y = Inches(1.5 + i * 1.8)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y, Inches(6), Inches(1.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = layer['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        title_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.1), Inches(2), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = layer['name']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        tech_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.6), Inches(5.4), Inches(0.8))
        tb = tech_box.text_frame
        for tech in layer['techs']:
            p = tb.add_paragraph()
            p.text = f'• {tech}'
            p.font.size = Pt(10)
            p.font.color.rgb = COLORS['ink']
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.5), Inches(7), Inches(4))
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['blue']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    stacks = [
        {'title': '前端技术栈', 'items': ['Vue 3.5.x', 'Vite 8.1.x', 'Element Plus 2.14.x', 'TypeScript 6.0.x']},
        {'title': '后端技术栈', 'items': ['Spring Boot 3.4.x', 'MyBatis Plus 3.5.x', 'MySQL 8.0', 'Redis 7.0']},
        {'title': '部署架构', 'items': ['Docker 24.x', 'Docker Compose 2.x', 'Nginx 1.25.x', 'TRTC 5.18.x']},
    ]
    
    stack_y = Inches(1.8)
    for stack in stacks:
        title_box = slide.shapes.add_textbox(Inches(8.3), stack_y, Inches(6.4), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = stack['title']
        tp.font.size = Pt(13)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        item_box = slide.shapes.add_textbox(Inches(8.3), stack_y + Inches(0.5), Inches(6.4), Inches(0.8))
        ib = item_box.text_frame
        for item in stack['items']:
            p = ib.add_paragraph()
            p.text = f'• {item}'
            p.font.size = Pt(10)
            p.font.color.rgb = COLORS['ink']
        
        stack_y += Inches(1.4)

def add_innovation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '09 / 12  技术创新'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    innovations = [
        {'icon': '🤖', 'title': 'AI驱动学习', 'desc': '深度整合大语言模型', 'color': COLORS['green']},
        {'icon': '💰', 'title': '学分银行体系', 'desc': '学习成果量化', 'color': COLORS['yellow']},
        {'icon': '📹', 'title': '实时视频面试', 'desc': '整合TRTC', 'color': COLORS['blue']},
        {'icon': '🎨', 'title': 'Neo-brutal风格', 'desc': '独特视觉风格', 'color': COLORS['pink']},
        {'icon': '⚡', 'title': '前后端分离', 'desc': 'Vue3+Spring Boot', 'color': COLORS['purple']},
    ]
    
    for i, item in enumerate(innovations):
        row = i // 3
        col = i % 3
        x = Inches(1 + col * 5)
        y = Inches(1.5 + row * 2.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(2.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = item['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(1.7), y + Inches(0.3), Inches(1), Inches(0.6))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = item['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1), Inches(3.5), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = item['title']
        tp.font.size = Pt(13)
        tp.font.color.rgb = COLORS['white'] if item['color'] == COLORS['green'] else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1.5), Inches(3.5), Inches(0.5))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = item['desc']
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['white'] if item['color'] == COLORS['green'] else COLORS['ink']
        dp.alignment = PP_ALIGN.CENTER
    
    footer_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(8), Inches(0.6))
    fb = footer_box.text_frame
    fp = fb.add_paragraph()
    fp.text = '5大创新点，打造差异化竞争力'
    fp.font.size = Pt(16)
    fp.font.color.rgb = COLORS['green']
    fp.font.bold = True
    fp.alignment = PP_ALIGN.CENTER

def add_process(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '10 / 12  流程演示'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    processes = [
        {'icon': '🎤', 'title': '面试流程', 'steps': ['浏览职位', '投递简历', '面试邀请', '在线面试', '面试结果'], 'color': COLORS['blue']},
        {'icon': '🎉', 'title': '活动流程', 'steps': ['发布活动', '邀请学员', '报名参加', '活动签到', '获取学分'], 'color': COLORS['pink']},
        {'icon': '📚', 'title': '学习流程', 'steps': ['选择课程', '开始学习', '互动交流', '完成课程', '积累学分'], 'color': COLORS['green']},
    ]
    
    for i, process in enumerate(processes):
        x = Inches(1 + i * 5)
        y = Inches(1.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(4))
        rect.fill.solid()
        rect.fill.fore_color.rgb = process['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(1.7), y + Inches(0.3), Inches(1), Inches(0.6))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = process['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1), Inches(3.5), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = process['title']
        tp.font.size = Pt(14)
        tp.font.color.rgb = COLORS['white'] if process['color'] == COLORS['green'] else COLORS['ink']
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        step_y = y + Inches(1.6)
        for j, step in enumerate(process['steps']):
            step_box = slide.shapes.add_textbox(x + Inches(0.3), step_y, Inches(4), Inches(0.5))
            sb = step_box.text_frame
            sp = sb.add_paragraph()
            sp.text = f'{j+1:02d}. {step}'
            sp.font.size = Pt(11)
            sp.font.color.rgb = COLORS['white'] if process['color'] == COLORS['green'] else COLORS['ink']
            step_y += Inches(0.55)
    
    footer_box = slide.shapes.add_textbox(Inches(4), Inches(6), Inches(8), Inches(0.6))
    fb = footer_box.text_frame
    fp = fb.add_paragraph()
    fp.text = '三大核心流程，连接学员、课程与企业'
    fp.font.size = Pt(16)
    fp.font.color.rgb = COLORS['green']
    fp.font.bold = True
    fp.alignment = PP_ALIGN.CENTER

def add_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['cream']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '11 / 12  组织管理'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['ink']
    p.font.bold = True
    
    members = [
        {'icon': '👨‍💻', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]', 'color': COLORS['blue']},
        {'icon': '👩‍💻', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]', 'color': COLORS['pink']},
        {'icon': '👨‍🎨', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]', 'color': COLORS['yellow']},
        {'icon': '👩‍🎨', 'name': '[姓名]', 'role': '[角色]', 'desc': '[职责描述]', 'color': COLORS['purple']},
    ]
    
    for i, member in enumerate(members):
        x = Inches(1 + i * 3.8)
        y = Inches(1.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.5), Inches(2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = member['color']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(3)
        
        icon_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(1), Inches(0.8))
        ib = icon_box.text_frame
        ip = ib.add_paragraph()
        ip.text = member['icon']
        ip.font.size = Pt(28)
        ip.alignment = PP_ALIGN.CENTER
        
        info_box = slide.shapes.add_textbox(x + Inches(1.5), y + Inches(0.2), Inches(1.7), Inches(1.5))
        ib = info_box.text_frame
        p = ib.add_paragraph()
        p.text = member['name']
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['ink']
        p.font.bold = True
        
        p = ib.add_paragraph()
        p.text = member['role']
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['ink']
        
        p = ib.add_paragraph()
        p.text = member['desc']
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['ink']
    
    workflow_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(14), Inches(2))
    workflow_rect.fill.solid()
    workflow_rect.fill.fore_color.rgb = COLORS['white']
    workflow_rect.line.color.rgb = COLORS['ink']
    workflow_rect.line.width = Pt(3)
    
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(4.1), Inches(5), Inches(0.4))
    tb = title_box.text_frame
    tp = tb.add_paragraph()
    tp.text = '团队协作流程'
    tp.font.size = Pt(14)
    tp.font.color.rgb = COLORS['ink']
    tp.font.bold = True
    
    workflows = [
        ('01', '需求分析', '收集用户需求'),
        ('02', '系统设计', '架构与数据库设计'),
        ('03', '开发实现', '前后端并行开发'),
        ('04', '部署上线', 'Docker容器化部署'),
    ]
    
    for i, (num, title, desc) in enumerate(workflows):
        x = Inches(1 + i * 3.5)
        y = Inches(4.6)
        
        num_box = slide.shapes.add_textbox(x, y, Inches(0.8), Inches(0.6))
        nb = num_box.text_frame
        np = nb.add_paragraph()
        np.text = num
        np.font.size = Pt(20)
        np.font.color.rgb = COLORS['green']
        np.font.bold = True
        
        title_box = slide.shapes.add_textbox(x + Inches(0.8), y, Inches(1.5), Inches(0.4))
        tb = title_box.text_frame
        tp = tb.add_paragraph()
        tp.text = title
        tp.font.size = Pt(12)
        tp.font.color.rgb = COLORS['ink']
        tp.font.bold = True
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.4), Inches(1.5), Inches(0.4))
        db = desc_box.text_frame
        dp = db.add_paragraph()
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = COLORS['ink']

def add_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS['green']
    
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(0.8))
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = '12 / 12  项目小结'
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True
    
    left_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(6.5), Inches(3.5))
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = COLORS['white']
    left_rect.line.color.rgb = COLORS['ink']
    left_rect.line.width = Pt(3)
    
    title_box1 = slide.shapes.add_textbox(Inches(1.3), Inches(1.7), Inches(6), Inches(0.5))
    tb = title_box1.text_frame
    tp = tb.add_paragraph()
    tp.text = '项目成果'
    tp.font.size = Pt(16)
    tp.font.color.rgb = COLORS['ink']
    tp.font.bold = True
    
    items_box1 = slide.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(6), Inches(2.5))
    ib = items_box1.text_frame
    items1 = [
        '完整的学习平台（课程学习、学分银行、社区论坛）',
        'AI智能助手（个性化学习推荐和智能答疑）',
        '企业对接功能（企业入驻、职位发布、在线面试）',
        'Neo-brutal设计（独特视觉风格，提升品牌识别度）',
    ]
    for item in items1:
        p = ib.add_paragraph()
        p.text = f'[OK] {item}'
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['ink']
    
    right_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.5), Inches(6.5), Inches(3.5))
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = COLORS['white']
    right_rect.line.color.rgb = COLORS['ink']
    right_rect.line.width = Pt(3)
    
    title_box2 = slide.shapes.add_textbox(Inches(8.8), Inches(1.7), Inches(6), Inches(0.5))
    tb = title_box2.text_frame
    tp = tb.add_paragraph()
    tp.text = '未来展望'
    tp.font.size = Pt(16)
    tp.font.color.rgb = COLORS['ink']
    tp.font.bold = True
    
    items_box2 = slide.shapes.add_textbox(Inches(8.8), Inches(2.3), Inches(6), Inches(2.5))
    ib = items_box2.text_frame
    items2 = [
        '深化AI应用（AI课程生成、智能评估）',
        '扩展社交功能（学习小组、好友互动）',
        '移动端适配（小程序和APP）',
        '数据分析平台（可视化报表和洞察）',
    ]
    for item in items2:
        p = ib.add_paragraph()
        p.text = f'[IDEA] {item}'
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['ink']
    
    numbers = [('6', '核心模块'), ('5', '技术创新'), ('3', '核心流程'), ('4', '团队成员')]
    for i, (num, label) in enumerate(numbers):
        x = Inches(1 + i * 3.8)
        y = Inches(5.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.5), Inches(1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS['white']
        rect.line.color.rgb = COLORS['ink']
        rect.line.width = Pt(2)
        
        num_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.1), Inches(1), Inches(0.8))
        nb = num_box.text_frame
        np = nb.add_paragraph()
        np.text = num
        np.font.size = Pt(28)
        np.font.color.rgb = COLORS['green']
        np.font.bold = True
        
        label_box = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.2), Inches(2), Inches(0.6))
        lb = label_box.text_frame
        lp = lb.add_paragraph()
        lp.text = label
        lp.font.size = Pt(12)
        lp.font.color.rgb = COLORS['ink']
        lp.font.bold = True
    
    closing_box = slide.shapes.add_textbox(Inches(5), Inches(7), Inches(6), Inches(1))
    cb = closing_box.text_frame
    cp = cb.add_paragraph()
    cp.text = '感谢聆听！'
    cp.font.size = Pt(28)
    cp.font.color.rgb = COLORS['white']
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER

if __name__ == '__main__':
    create_presentation()